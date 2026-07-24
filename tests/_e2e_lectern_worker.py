#!/usr/bin/env python3
"""Lectern-side worker for the synthetic cross-repository golden path."""

from __future__ import annotations

import argparse
import json
import sys
from pathlib import Path
from typing import Any

from asa_agents import FakeLLM, build_outline
from asa_agents.graph import build_graph
from ingestion import ingest_handoff
from langgraph.types import Command
from pptx import Presentation
from slide_ir import GenerationState


class WorkerFailure(RuntimeError):
    """Lectern violated a golden-path invariant."""


def require(condition: bool, message: str) -> None:
    if not condition:
        raise WorkerFailure(message)


def read_json(path: Path) -> dict[str, Any]:
    value = json.loads(path.read_text(encoding="utf-8"))
    require(isinstance(value, dict), "Handoff metadata must be an object")
    return value


def scripted_deck(meta: dict[str, Any]) -> str:
    papers = meta["papers"]
    rows = [
        [
            str(paper["title"]),
            f"SYNTHETIC-PAPER-{index}-EVIDENCE",
            str(paper["doi"]),
        ]
        for index, paper in enumerate(papers, start=1)
    ]
    return json.dumps(
        {
            "deck_id": "synthetic-lectern-golden",
            "slides": [
                {
                    "slide_id": "cover",
                    "layout_type": "title",
                    "title": "[SYNTHETIC] Dual-paper evidence review",
                    "subtitle": "Offline, credential-free golden path",
                    "blocks": [],
                    "speaker_notes": "This deck contains synthetic evidence only.",
                    "provenance": {"source": "report_basis"},
                },
                {
                    "slide_id": "evidence",
                    "layout_type": "two_column_table",
                    "title": "Evidence-backed comparison",
                    "blocks": [
                        {
                            "type": "table",
                            "columns": ["Source", "Evidence marker", "DOI"],
                            "rows": rows,
                            "needs_human_check": False,
                        }
                    ],
                    "speaker_notes": "Each row remains traceable to the handoff.",
                    "provenance": {
                        "source": f"handoff:{meta['key']}",
                        "paper_count": len(papers),
                    },
                },
            ],
        },
        ensure_ascii=True,
    )


def all_text(prs: Presentation) -> list[str]:
    values: list[str] = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if getattr(shape, "has_text_frame", False):
                values.append(shape.text)
            if getattr(shape, "has_table", False):
                table = shape.table
                for row in table.rows:
                    values.extend(cell.text for cell in row.cells)
    return values


def field(value: Any, name: str, default: Any = None) -> Any:
    if isinstance(value, dict):
        return value.get(name, default)
    return getattr(value, name, default)


def run(handoff: Path, output_dir: Path) -> dict[str, Any]:
    meta = read_json(handoff / "meta.json")
    papers = meta.get("papers")
    require(meta.get("schema_version") == "handoff/1.1", "Expected handoff/1.1")
    require(isinstance(papers, list) and len(papers) == 2, "Expected two papers")

    output_dir.mkdir(parents=True, exist_ok=True)
    require(not list(output_dir.rglob("*.pptx")), "Output directory was not empty")

    ingested = ingest_handoff(handoff)
    evidence_text = "\n".join(
        str(asset.content_ref)
        for asset in ingested.assets
        if str(asset.kind.value) == "section_text"
    )
    expected_evidence = [
        str(paper[field])
        for paper in papers
        for field in ("title", "doi")
    ] + [
        "SYNTHETIC-PAPER-1-EVIDENCE",
        "SYNTHETIC-PAPER-2-EVIDENCE",
    ]
    for value in expected_evidence:
        require(value in evidence_text, f"Lectern evidence omitted {value!r}")

    llm = FakeLLM(scripted_deck(meta))
    state = GenerationState(
        job_id="synthetic-golden-path",
        evidence=ingested.assets,
        tables=ingested.tables,
        max_retries=0,
    )
    graph = build_graph(
        llm,
        out_dir=output_dir,
        planner=build_outline,
    )
    config = {"configurable": {"thread_id": state.job_id}}
    graph.invoke(state.model_dump(), config)
    snapshot = graph.get_state(config)

    require(
        "approval" in snapshot.next,
        "Production graph did not pause at the outline approval gate",
    )
    require(
        not field(snapshot.values, "user_approved_outline", False),
        "Production graph bypassed human approval",
    )
    require(
        field(snapshot.values, "output_path") is None,
        "Production graph assigned an output before approval",
    )
    require(
        not list(output_dir.rglob("*.pptx")),
        "A persistent PPTX was created before outline approval",
    )
    require(len(llm.calls) == 1, "FakeLLM was not called exactly once")
    prompt = llm.calls[0]["prompt"]
    for value in expected_evidence:
        require(value in prompt, f"Planner prompt omitted {value!r}")

    final = graph.invoke(Command(resume={"approved": True}), config)
    require(field(final, "user_approved_outline") is True, "Approval was not recorded")
    approved_value = field(final, "output_path")
    require(isinstance(approved_value, str), "Approval did not produce an output path")
    approved_path = Path(approved_value)
    require(approved_path.is_file(), "Approved deck was not generated")

    presentation = Presentation(str(approved_path))
    require(len(presentation.slides) == 2, "Generated deck has the wrong slide count")
    table_shapes = [
        shape
        for slide in presentation.slides
        for shape in slide.shapes
        if getattr(shape, "has_table", False)
    ]
    require(len(table_shapes) == 1, "Generated deck lacks one native table")
    table = table_shapes[0].table
    require(len(table.rows) == 3, "Native table lost a synthetic evidence row")
    require(len(table.columns) == 3, "Native table lost a column")

    title_shape = next(
        (
            shape
            for slide in presentation.slides
            for shape in slide.shapes
            if getattr(shape, "has_text_frame", False)
            and "[SYNTHETIC] Dual-paper evidence review" in shape.text
        ),
        None,
    )
    require(title_shape is not None, "Generated deck lacks editable title text")
    title_shape.text_frame.text = "[SYNTHETIC] Edited evidence review"
    table.cell(1, 1).text = "SYNTHETIC-EDIT-CONFIRMED"

    edited_path = output_dir / "edited-synthetic-deck.pptx"
    presentation.save(str(edited_path))
    reopened = Presentation(str(edited_path))
    reopened_text = all_text(reopened)
    require(
        "[SYNTHETIC] Edited evidence review" in reopened_text,
        "Edited title did not survive reopening",
    )
    require(
        "SYNTHETIC-EDIT-CONFIRMED" in reopened_text,
        "Edited table cell did not survive reopening",
    )

    return {
        "status": "passed",
        "handoff_schema": meta["schema_version"],
        "paper_count": len(papers),
        "evidence_asset_count": len(ingested.assets),
        "warning_count": len(ingested.warnings),
        "production_graph_exercised": True,
        "planner_used_source_evidence": True,
        "hard_stop_blocked_pptx": True,
        "approved_slide_count": len(reopened.slides),
        "native_table_count": len(table_shapes),
        "editable_after_reopen": True,
    }


def build_parser() -> argparse.ArgumentParser:
    parser = argparse.ArgumentParser(description=__doc__)
    parser.add_argument("--handoff", type=Path, required=True)
    parser.add_argument("--output-dir", type=Path, required=True)
    return parser


def main(argv: list[str] | None = None) -> int:
    args = build_parser().parse_args(argv)
    try:
        report = run(args.handoff.resolve(), args.output_dir.resolve())
    except (OSError, ValueError, WorkerFailure) as exc:
        print(f"FAIL: {exc}", file=sys.stderr)
        return 1
    print(json.dumps(report, ensure_ascii=True, sort_keys=True))
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
